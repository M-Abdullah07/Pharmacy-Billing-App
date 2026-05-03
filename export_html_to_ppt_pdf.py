# Dependencies:
# pip install Pillow python-pptx playwright
# playwright install msedge

import sys
import os
import time
import threading
import http.server
import socketserver
import argparse
import shutil
from pathlib import Path
from PIL import Image, JpegImagePlugin, PdfImagePlugin
from pptx import Presentation
from pptx.util import Inches
from playwright.sync_api import sync_playwright

# Native MS Edge approach avoids bundling Chromium binaries
def find_free_port():
    """Finds an available port on the system."""
    with socketserver.TCPServer(("", 0), None) as s:
        return s.socket.getsockname()[1]

def start_server(directory, port):
    """Starts a simple HTTP server in a separate thread."""
    class Handler(http.server.SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, directory=directory, **kwargs)
        
        def log_message(self, format, *args):
            # Silence server logs for a cleaner CLI experience
            pass

    try:
        with socketserver.TCPServer(("", port), Handler) as httpd:
            httpd.serve_forever()
    except Exception as e:
        print(f"Server error: {e}")

def export_presentation(args):
    """Main export logic."""
    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"Error: Input file {input_path} does not exist.")
        sys.exit(1)

    html_file = input_path.name
    serve_dir = str(input_path.parent)
    
    # Configure output paths
    base_name = input_path.stem
    output_pdf = args.pdf_output or str(input_path.with_suffix(".pdf"))
    output_pptx = args.pptx_output or str(input_path.with_suffix(".pptx"))
    screenshot_dir = Path(args.screenshot_dir)
    
    # Ensure screenshot directory exists
    screenshot_dir.mkdir(parents=True, exist_ok=True)
    
    # Start server
    port = args.port if args.port else find_free_port()
    server_thread = threading.Thread(target=start_server, args=(serve_dir, port), daemon=True)
    server_thread.start()
    time.sleep(1.5)  # Give server time to spin up

    screenshot_paths = []

    print(f"Starting browser (Viewport: {args.width}x{args.height})...")
    with sync_playwright() as p:
        # Using native Edge browser eliminates the need to bundle Chromium binaries
        browser = p.chromium.launch(channel="msedge", headless=True)
        page = browser.new_page(viewport={"width": args.width, "height": args.height})
        
        url = f"http://localhost:{port}/{html_file}"
        print(f"Loading {url}...")
        page.goto(url, wait_until="networkidle")
        
        # Wait for fonts and initial animations
        page.evaluate("document.fonts.ready")
        time.sleep(args.initial_delay)

        # Detect slides
        slide_count = page.evaluate("document.querySelectorAll('.slide').length")
        if slide_count == 0:
            print("Error: No elements with class '.slide' found.")
            browser.close()
            return

        print(f"Found {slide_count} slides. Capturing screenshots...")

        for i in range(slide_count):
            # Force visibility and reset state for the target slide
            page.evaluate(f"""(index) => {{
                const slides = document.querySelectorAll('.slide');
                slides.forEach((slide, idx) => {{
                    if (idx === index) {{
                        slide.style.display = 'flex';
                        slide.style.opacity = '1';
                        slide.style.visibility = 'visible';
                        slide.style.position = 'relative';
                        slide.style.transform = 'none';
                        slide.classList.add('active');
                        slide.classList.add('visible');
                    }} else {{
                        slide.style.display = 'none';
                        slide.classList.remove('active');
                        slide.classList.remove('visible');
                    }}
                }});
                
                // Force reveal animations
                const currentSlide = slides[index];
                if (currentSlide) {{
                    currentSlide.querySelectorAll('.reveal').forEach(el => {{
                        el.style.opacity = '1';
                        el.style.transform = 'none';
                        el.style.visibility = 'visible';
                    }});
                }}
            }}""", i)
            
            # Brief pause for any CSS transitions triggered by visibility
            time.sleep(args.slide_delay)
            
            path = screenshot_dir / f"slide_{i:03d}.png"
            page.screenshot(path=str(path))
            screenshot_paths.append(path)
            print(f"[{i+1}/{slide_count}] Captured slide")
            
        browser.close()

    if screenshot_paths:
        # PDF Generation
        if not args.no_pdf:
            print(f"Generating PDF: {output_pdf}")
            images = [Image.open(p).convert("RGB") for p in screenshot_paths]
            images[0].save(output_pdf, save_all=True, append_images=images[1:], format="PDF")
            print("Success: PDF Exported")

        # PPTX Generation
        if not args.no_pptx:
            print(f"Generating PPTX: {output_pptx}")
            prs = Presentation()
            # Standard 16:9 Slide size
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6] # Blank slide index
            
            for p in screenshot_paths:
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            prs.save(output_pptx)
            print("Success: PPTX Exported")


    # Cleanup
    if not args.keep_screenshots:
        print("Cleaning up temporary screenshots...")
        shutil.rmtree(screenshot_dir)

def main():
    parser = argparse.ArgumentParser(description="Professional HTML Presentation to PDF/PPTX Converter")
    
    # Required
    parser.add_argument("input", help="Path to the input HTML file")
    
    # Optional outputs
    parser.add_argument("--pdf-output", help="Custom path for the PDF output")
    parser.add_argument("--pptx-output", help="Custom path for the PPTX output")
    parser.add_argument("--no-pdf", action="store_true", help="Disable PDF generation")
    parser.add_argument("--no-pptx", action="store_true", help="Disable PPTX generation")
    
    # Display Settings
    parser.add_argument("--width", type=int, default=1920, help="Viewport width (default: 1920)")
    parser.add_argument("--height", type=int, default=1080, help="Viewport height (default: 1080)")
    
    # Timing
    parser.add_argument("--initial-delay", type=float, default=2.0, help="Wait time for fonts/assets (seconds)")
    parser.add_argument("--slide-delay", type=float, default=0.5, help="Wait time per slide (seconds)")
    
    # Infrastructure
    parser.add_argument("--port", type=int, help="Port for the temporary server (defaults to a free port)")
    parser.add_argument("--screenshot-dir", default="temp_screenshots", help="Temporary directory for screenshots")
    parser.add_argument("--keep-screenshots", action="store_true", help="Do not delete the screenshot directory after completion")

    args = parser.parse_args()
    
    try:
        export_presentation(args)
    except KeyboardInterrupt:
        print("\nExport cancelled by user.")
        sys.exit(0)
    except Exception as e:
        print(f"\nCritical Error: {e}")
        sys.exit(1)

if __name__ == '__main__':
    main()

